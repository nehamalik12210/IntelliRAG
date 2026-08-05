"""Multi-format document loaders.

Supports: PDF, DOCX, TXT, Markdown, CSV, PPTX, HTML
Each loader returns a list of dicts with 'content' and 'metadata'.
"""

import csv
import io
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def load_document(file_path: str, file_type: str) -> list[dict]:
    """Load a document and return a list of page/section dicts.

    Each dict contains:
        - content: str — the text content
        - metadata: dict — {page_number, source_filename, ...}

    Args:
        file_path: Path to the uploaded file
        file_type: File extension (e.g., 'pdf', 'docx', 'txt')

    Returns:
        List of page/section dicts
    """
    file_type = file_type.lower().strip(".")
    filename = Path(file_path).name

    loaders = {
        "pdf": _load_pdf,
        "docx": _load_docx,
        "txt": _load_text,
        "md": _load_text,
        "csv": _load_csv,
        "pptx": _load_pptx,
        "html": _load_html,
        "htm": _load_html,
    }

    loader = loaders.get(file_type)
    if not loader:
        raise ValueError(f"Unsupported file type: {file_type}")

    pages = loader(file_path)

    # Attach filename to all metadata
    for page in pages:
        page["metadata"]["source_filename"] = filename

    logger.info(f"Loaded {len(pages)} pages from '{filename}' ({file_type})")
    return pages


def _load_pdf(file_path: str) -> list[dict]:
    """Load PDF using PyPDF2, one entry per page."""
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({
                "content": text,
                "metadata": {"page_number": i + 1},
            })
    return pages


def _load_docx(file_path: str) -> list[dict]:
    """Load DOCX using python-docx, concatenate all paragraphs."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []

    # Treat whole document as single page (DOCX doesn't have page concepts)
    return [{
        "content": "\n\n".join(paragraphs),
        "metadata": {"page_number": 1},
    }]


def _load_text(file_path: str) -> list[dict]:
    """Load plain text or markdown files."""
    import chardet

    with open(file_path, "rb") as f:
        raw = f.read()

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = raw.decode(encoding, errors="replace")

    if not text.strip():
        return []

    return [{
        "content": text,
        "metadata": {"page_number": 1},
    }]


def _load_csv(file_path: str) -> list[dict]:
    """Load CSV, converting each row to a readable text format."""
    import chardet

    with open(file_path, "rb") as f:
        raw = f.read()

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    text = raw.decode(encoding, errors="replace")

    reader = csv.DictReader(io.StringIO(text))
    rows = []
    for i, row in enumerate(reader):
        row_text = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
        if row_text.strip():
            rows.append(row_text)

    if not rows:
        return []

    # Group rows into pages of 50 for chunking
    page_size = 50
    pages = []
    for i in range(0, len(rows), page_size):
        page_rows = rows[i:i + page_size]
        pages.append({
            "content": "\n".join(page_rows),
            "metadata": {"page_number": (i // page_size) + 1},
        })
    return pages


def _load_pptx(file_path: str) -> list[dict]:
    """Load PowerPoint, one entry per slide."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append({
                "content": "\n".join(texts),
                "metadata": {"page_number": i + 1, "slide_number": i + 1},
            })
    return slides


def _load_html(file_path: str) -> list[dict]:
    """Load HTML by extracting text with BeautifulSoup."""
    from bs4 import BeautifulSoup
    import chardet

    with open(file_path, "rb") as f:
        raw = f.read()

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"
    html = raw.decode(encoding, errors="replace")

    # Use BeautifulSoup to properly strip script/style tags and decode entities
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator=" ", strip=True)

    if not text:
        return []

    return [{
        "content": text,
        "metadata": {"page_number": 1},
    }]
